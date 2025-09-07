from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
from pptx import Presentation
from docx import Document
import os


def extract_pdf_pages(path: str):
    pages = []
    for page_layout in extract_pages(path):
        page_text = []
        for element in page_layout:
            if isinstance(element, LTTextContainer):
                page_text.append(element.get_text())
        pages.append("\n".join(page_text))
    return pages


def extract_pptx_slides(path: str):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append("\n".join(texts))
    return slides


def extract_docx_paragraphs(path: str):
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return paragraphs


def extract_text_by_type(path: str, content_type: str):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_pages(path)
    if ext in (".pptx", ".ppt"):
        return extract_pptx_slides(path)
    if ext in (".docx", ".doc"):
        return ["\n\n".join(extract_docx_paragraphs(path))]
    # fallback: read as text
    try:
        with open(path, "r", encoding="utf-8") as f:
            return [f.read()]
    except Exception:
        return [""]


