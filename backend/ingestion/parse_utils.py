import os


def extract_pdf_pages(path: str):
    # Lazy import heavy deps to avoid import-time failures in environments
    # where pdf dependencies are not installed unless actually used.
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer, LAParams

    laparams = LAParams(
        line_margin=0.15,   # tighter line grouping
        char_margin=2.0,    # more tolerant character grouping
        word_margin=0.1,    # conservative word spacing
        boxes_flow=None,    # disable automatic layout flow assumptions
        all_texts=True,
    )
    pages = []
    for page_layout in extract_pages(path, laparams=laparams):
        page_text = []
        for element in page_layout:
            if isinstance(element, LTTextContainer):
                txt = element.get_text()
                # minor cleanup: drop extremely short fragments
                if txt and len(txt.strip()) >= 2:
                    page_text.append(txt)
        pages.append("\n".join(page_text))
    return pages


def extract_pptx_slides(path: str):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append("\n".join(texts))
    return slides


def extract_docx_paragraphs(path: str):
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return paragraphs


def extract_text_by_type(path: str, content_type: str):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_pages(path)
    if ext in (".pptx", ".ppt"):
        return extract_pptx_slides(path)
    if ext in (".docx", ".doc"):
        return ["\n\n".join(extract_docx_paragraphs(path))]
    # fallback: read as text
    try:
        with open(path, "r", encoding="utf-8") as f:
            return [f.read()]
    except Exception:
        return [""]


